import os
from pypdf import PdfReader
from docx import Document
from pptx import Presentation


def read_pdf(file_path):
    """Extract text from a PDF file"""
    text = []
    reader = PdfReader(file_path)

    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text.append(page_text)

    return "\n".join(text)


def read_docx(file_path):
    """Extract text from a Word document"""
    doc = Document(file_path)
    return "\n".join(paragraph.text for paragraph in doc.paragraphs)


def read_pptx(file_path):
    """Extract text from a PowerPoint file"""
    prs = Presentation(file_path)
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return "\n".join(text)


def load_all_documents(upload_dir="data/uploads"):
    """
    Reads ALL uploaded files (PDF, DOCX, PPTX)
    and returns combined text
    """
    all_text = []

    for root, _, files in os.walk(upload_dir):
        for file in files:
            file_path = os.path.join(root, file)

            if file.lower().endswith(".pdf"):
                all_text.append(read_pdf(file_path))

            elif file.lower().endswith(".docx"):
                all_text.append(read_docx(file_path))

            elif file.lower().endswith(".pptx"):
                all_text.append(read_pptx(file_path))

    return "\n".join(all_text)
